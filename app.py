from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
from youtube_transcript_api import YouTubeTranscriptApi
import google.generativeai as genai
import tempfile
import os
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__, template_folder='SEREN AI/templates')
CORS(app)  # Enable CORS for all routes

# Configuration
app.config.update({
    'YOUTUBE_API_KEY': os.getenv('YOUTUBE_API_KEY'),
    'GEMINI_API_KEY': os.getenv('GEMINI_API_KEY'),
    'SECRET_KEY': os.getenv('SECRET_KEY'),
    'UPLOAD_FOLDER': tempfile.mkdtemp(),
    'MAX_CONTENT_LENGTH': 16 * 1024 * 1024,  # 16MB limit
    'ALLOWED_EXTENSIONS': {'pdf', 'ppt', 'pptx', 'txt', 'doc', 'docx'}
})

# Initialize Gemini
genai.configure(api_key=app.config['GEMINI_API_KEY'])

# Error handlers
@app.errorhandler(400)
def bad_request(error):
    return jsonify({'error': 'Bad request', 'message': str(error)}), 400

@app.errorhandler(404)
def not_found(error):
    return jsonify({'error': 'Not found', 'message': str(error)}), 404

@app.errorhandler(500)
def internal_error(error):
    return jsonify({'error': 'Internal server error', 'message': str(error)}), 500

# Utility functions
def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def extract_text_from_file(filepath):
    try:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.pdf':
            with open(filepath, 'rb') as f:
                return "\n".join(page.extract_text() for page in PdfReader(f).pages if page.extract_text())
        elif ext in ('.ppt', '.pptx'):
            return "\n".join(shape.text for slide in Presentation(filepath).slides 
                          for shape in slide.shapes if hasattr(shape, "text") and shape.text)
        elif ext in ('.doc', '.docx'):
            return "\n".join(p.text for p in docx.Document(filepath).paragraphs if p.text)
        elif ext == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        raise ValueError(f"Unsupported file format: {ext}")
    except Exception as e:
        logger.error(f"Error extracting text: {str(e)}")
        raise ValueError(f"Error processing file: {str(e)}")

def generate_quiz(text, num_questions=5):
    try:
        prompt = f"""Generate {num_questions} clear multiple-choice questions based on the following text.
For each question, provide 4 distinct options (A-D) and indicate the correct answer.

Text:
{text}

Format each question exactly like this:
Question: [question text]
A) Option 1
B) Option 2
C) Option 3
D) Option 4
Answer: [correct letter]

Ensure questions cover key concepts from the text."""
        
        model = genai.GenerativeModel('gemini-1.5-pro-latest')
        response = model.generate_content(prompt)
        
        if not response.text:
            raise ValueError("No quiz generated - empty response from AI")
            
        return response.text
    except Exception as e:
        logger.error(f"Quiz generation failed: {str(e)}")
        raise ValueError(f"Failed to generate quiz: {str(e)}")

# Routes
@app.route('/')
def home():
    return render_template('Index.html')

@app.route('/dashboard')
def dashboard():
    return render_template('dashboard.html')

@app.route('/notes')
def notes():
    return render_template('notes.html', api_key=app.config['GEMINI_API_KEY'])

@app.route('/quizzes')
def quizzes():
    return render_template('quizzes.html')

@app.route('/saved_videos')
def saved_videos():
    return render_template('saved_videos.html')

@app.route('/sign')
def sign():
    return render_template('sign.html')

@app.route('/profile')
def user_profile():
    return render_template('user_profile.html')

@app.route('/video-recommendations')
def video_reco():
    return render_template('video_reco.html')

@app.route('/study_planner')
def study_planner():
    return render_template('study_planner.html')

@app.route('/get_transcript', methods=['GET'])
def get_transcript():
    video_id = request.args.get('video_id')
    if not video_id:
        return jsonify({"error": "video_id parameter is required"}), 400
    
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        transcript_text = " ".join([entry['text'] for entry in transcript])
        return jsonify({"transcript": transcript_text})
    except Exception as e:
        logger.error(f"Transcript error for video {video_id}: {str(e)}")
        return jsonify({"error": f"Failed to get transcript: {str(e)}"}), 500

@app.route('/generate_summary', methods=['POST'])
def generate_summary():
    if not request.is_json:
        return jsonify({"error": "Request must be JSON"}), 400
        
    data = request.get_json()
    transcript = data.get('transcript', '')
    
    if not transcript:
        return jsonify({"error": "No transcript provided"}), 400
    
    try:
        model = genai.GenerativeModel('gemini-1.5-pro-latest')
        prompt = f"""Generate a concise, well-structured summary (about 100 words) of the following video transcript.
Focus on key points and main ideas. Use bullet points if appropriate.

Transcript:
{transcript}"""
        response = model.generate_content(prompt)
        
        if not response.text:
            return jsonify({"error": "Failed to generate summary"}), 500
            
        return jsonify({"summary": response.text})
    except Exception as e:
        logger.error(f"Summary generation failed: {str(e)}")
        return jsonify({"error": f"Summary generation failed: {str(e)}"}), 500

@app.route('/api/generate-quiz-from-file', methods=['POST'])
def handle_file_upload():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({
            'error': 'File type not allowed',
            'allowed_types': list(app.config['ALLOWED_EXTENSIONS'])
        }), 400

    try:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        text = extract_text_from_file(filepath)
        if not text.strip():
            raise ValueError("Extracted text is empty - file may be corrupted or contain no text")
            
        quiz = generate_quiz(text)
        
        return jsonify({
            'quiz': quiz,
            'filename': filename,
            'text_length': len(text),
            'num_questions': quiz.count('Question:')
        })
    except Exception as e:
        logger.error(f"File upload error: {str(e)}")
        return jsonify({'error': str(e)}), 500
    finally:
        if 'filepath' in locals() and os.path.exists(filepath):
            os.remove(filepath)

@app.route('/api/generate-quiz-from-text', methods=['POST'])
def generate_quiz_from_text():
    if not request.is_json:
        return jsonify({'error': 'Request must be JSON'}), 400
    
    data = request.get_json()
    text = data.get('text', '')
    num_questions = min(int(data.get('num_questions', 5)), 10)  # Limit to 10 questions max
    
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    try:
        quiz = generate_quiz(text, num_questions)
        return jsonify({
            'quiz': quiz,
            'text_length': len(text),
            'num_questions': quiz.count('Question:')
        })
    except Exception as e:
        logger.error(f"Quiz from text error: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/chat', methods=['POST'])
def handle_chat():
    if not request.is_json:
        return jsonify({'error': 'Request must be JSON'}), 400
    
    data = request.get_json()
    message = data.get('message', '').strip()
    history = data.get('history', [])
    
    if not message:
        return jsonify({'error': 'No message provided'}), 400
    
    try:
        model = genai.GenerativeModel('gemini-1.5-pro-latest')
        
        # Build conversation context
        chat = model.start_chat(history=[])
        if history:
            for msg in history[-6:]:  # Keep last 6 messages for context
                if msg['role'] == 'user':
                    chat.send_message(msg['content'])
                else:
                    chat.history.append({'role': 'model', 'parts': [msg['content']]})
        
        response = chat.send_message(
            f"""As SEREN AI, a knowledgeable study assistant, provide helpful, accurate responses to student queries.
Current query: {message}
- Be concise but thorough
- Use bullet points when listing items
- Suggest related topics when appropriate
- Format answers clearly for readability""",
            generation_config={
                "max_output_tokens": 1000,
                "temperature": 0.3
            }
        )
        
        if not response.text:
            raise ValueError("Empty response from AI")
            
        return jsonify({
            'response': response.text,
            'status': 'success'
        })
    except Exception as e:
        logger.error(f"Chat error: {str(e)}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=os.getenv('FLASK_DEBUG', False))
